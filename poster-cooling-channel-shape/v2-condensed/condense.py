"""Condense poster text to key phrases, enlarge fonts, add two figures, re-flow.

Works on the author's edited poster (Times New Roman, 2-column A0 portrait),
preserving all styling, images, chart and colours.

Figures added:
  * current-density map (ICEM Fig. 3)  - beside the AC-loss chart, half width each
  * temperature fields  (ECCE Fig. 9)  - under the optimal-design table
Traded for the space: the "18 x" tile (it restated the table above it) and a
smaller Pareto plot.
"""
import re
from copy import deepcopy

from PIL import Image, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
SRC = "poster_v2.pptx"
OUT = "HiECSs_08_2026_poster_condensed.pptx"
FIGDIR = "/home/user/Claude/poster-cooling-channel-shape/figs/"

SERIF = "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf"
SERIF_B = "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf"
SERIF_BI = "/usr/share/fonts/truetype/liberation/LiberationSerif-BoldItalic.ttf"

BODY = 44          # was 32
BODY_NARROW = 44   # column beside the photo
HEAD = 52          # was 40
TABLE = 36         # was 32
QUESTION = 42      # was 35
CAPTION = 26
CHART_FONT = 22   # legend + axes, sized for the half-width chart
LNSPC = 1.12
FACE = 1.20        # Times New Roman line box vs. point size
SPC_AFT = 10       # pt after each bullet
MARL = Inches(0.34)
BULLET_PCT = 125000
GREY = RGBColor(0x59, 0x59, 0x59)

prs = Presentation(SRC)
slide = prs.slides[0]
sh = list(slide.shapes)


# ---------------------------------------------------------------- measuring
def n_lines(text, size_pt, width_in, bold=False, italic=False):
    path = SERIF_BI if (bold and italic) else SERIF_B if bold else SERIF
    font = ImageFont.truetype(path, size_pt)
    limit = width_in * 72.0
    lines, cur = 0, ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if font.getlength(trial) <= limit or not cur:
            cur = trial
        else:
            lines += 1
            cur = word
    return max(1, lines + (1 if cur else 0))


def text_height(items, size_pt, width_in, bold=False, italic=False, bullets=True):
    """Height in inches of a bulleted/plain paragraph list."""
    usable = width_in - (MARL.inches if bullets else 0)
    total = 0.0
    for i, t in enumerate(items):
        total += n_lines(t, size_pt, usable, bold, italic) * size_pt * LNSPC * FACE
        if i < len(items) - 1:
            total += SPC_AFT
    return total / 72.0


# ---------------------------------------------------------------- text edit
def set_items(shape, items, size_pt, bullets=True, algn="l"):
    """Replace a text frame's paragraphs, cloning the first one's formatting."""
    body = shape.text_frame._txBody
    paras = body.findall(A + "p")
    tpl = deepcopy(paras[0])
    for extra in tpl.findall(A + "r")[1:]:
        tpl.remove(extra)
    if not tpl.findall(A + "r"):
        raise ValueError("template paragraph has no run")
    for p in paras:
        body.remove(p)

    for text in items:
        p = deepcopy(tpl)
        pPr = p.find(A + "pPr")
        if pPr is None:
            pPr = p.makeelement(A + "pPr", {})
            p.insert(0, pPr)
        pPr.set("algn", algn)
        if bullets:
            pPr.set("marL", str(MARL))
            pPr.set("indent", str(-MARL))
            sz = pPr.find(A + "buSzPct")
            if sz is not None:
                sz.set("val", str(BULLET_PCT))
        run = p.find(A + "r")
        run.find(A + "t").text = text
        for rPr in (run.find(A + "rPr"), p.find(A + "endParaRPr")):
            if rPr is not None:
                rPr.set("sz", str(int(size_pt * 100)))
        body.append(p)


def set_table_font(table, size_pt, row_h_in):
    for row in table.rows:
        row.height = Inches(row_h_in)
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(size_pt)


def set_cell(table, r, c, text):
    cell = table.cell(r, c)
    runs = [run for p in cell.text_frame.paragraphs for run in p.runs]
    runs[0].text = text
    for extra in runs[1:]:
        extra.text = ""


def add_caption(text, width_in):
    tb = slide.shapes.add_textbox(0, 0, Inches(width_in), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name, run.font.size = "Times New Roman", Pt(CAPTION)
    run.font.italic, run.font.color.rgb = True, GREY
    tb.height = Inches(text_height([text], CAPTION, width_in, italic=True,
                                   bullets=False) + 0.05)
    return tb


def add_pic(path, width_in):
    ar = (lambda im: im.height / im.width)(Image.open(path))
    return slide.shapes.add_picture(path, 0, 0, Inches(width_in),
                                    Inches(round(width_in * ar, 4)))


# ---------------------------------------------------------------- content
MOTIVATION = [
    "Traction motors: higher speed and frequency",
    "Hairpin windings = industry standard",
    "AC losses → severe local hot spots",
    "Water jacket: long thermal path",
]
DOC = [
    "Oil flows inside each conductor",
    "Heat extracted at the source",
    "Hot spots suppressed",
    "+277 % continuous current",
]
EM = [
    "2D FEA: 350 A rms, 266.7 Hz",
    "Current crowds at the periphery",
    "Channel sits in the “dead zone”",
]
OPT = [
    "Bigger channel: more surface, less copper",
    "Smaller channel: steep pressure drop",
    "Temperature vs. pumping power compete",
    "128 FE samples → surrogate → GA",
    "Wall and opening ≥ 0.75 mm",
]
CONCL = [
    "Channel centre is EM “dead”: < 0.25 %",
    "Shape it purely for cooling",
    "Rectangle > ellipse ≫ circle",
    "Round → rectangular: −10 °C peak",
    "Optimum: 78.8 °C at 1.92 W",
    "Next: manufacturing + motorette test",
]
CAP_CURRENT = "(a) solid   (b) round   (c) rectangular"
CAP_CHART = "AC loss per layer (W)"
CAP_TEMP = "Optimal designs: (a) circular   (b) rectangular   (c) elliptical"

HEADERS = {6: "Motivation", 8: "Direct oil cooling", 10: "Reference machine",
           14: "Electromagnetic analysis", 17: "Channel-shape optimization",
           24: "Conclusions"}

# --------------------------------------- drop the "18 x" tile (restates table)
for i in (21, 22, 23):
    sh[i]._element.getparent().remove(sh[i]._element)

# ------------------------------------------------- normalize column geometry
COL1_X, COL1_W = 2.23, 13.32          # left column
COL2_X, COL2_W = 17.24, 13.64         # right column, right edge 30.88 = title


def set_box(idx, x, w):
    sh[idx].left, sh[idx].width = Inches(round(x, 4)), Inches(round(w, 4))


def set_table_box(idx, x, w):
    """Scale a table's columns so it spans exactly the column width."""
    tbl = sh[idx].table
    factor = Inches(w) / sum(c.width for c in tbl.columns)
    for col in tbl.columns:
        col.width = int(round(col.width * factor))
    set_box(idx, x, w)


def fit_picture(idx, x, w):
    """Resize a picture to a given width, preserving aspect ratio."""
    ar = Emu(sh[idx].height).inches / Emu(sh[idx].width).inches
    sh[idx].left, sh[idx].width = Inches(round(x, 4)), Inches(round(w, 4))
    sh[idx].height = Inches(round(w * ar, 4))


for i in (6, 7, 8, 10, 14, 15, 12):
    set_box(i, COL1_X, COL1_W)
for i in (17, 18, 24, 25):
    set_box(i, COL2_X, COL2_W)
set_box(13, COL1_X + 0.20, COL1_W - 0.40)
set_table_box(11, COL1_X, COL1_W)
set_table_box(20, COL2_X, COL2_W)
fit_picture(29, COL2_X, COL2_W)

# Pareto plot trimmed to make room for the temperature fields, kept centred
PARETO_H = 5.5
pareto_w = PARETO_H * (Emu(sh[19].width).inches / Emu(sh[19].height).inches)
fit_picture(19, COL2_X + (COL2_W - pareto_w) / 2, pareto_w)

# photo + text side by side: size the photo to the text, keep its aspect ratio
COL_L, COL_R_EDGE = COL1_X, COL1_X + COL1_W
PHOTO_AR = Emu(sh[28].width).inches / Emu(sh[28].height).inches
h9 = text_height(DOC, BODY_NARROW, 7.10)
photo_h = min(max(h9, 7.0), 8.6)
photo_w = photo_h * PHOTO_AR
x9 = COL_L + photo_w + 0.60
W9 = COL_R_EDGE - x9
h9 = text_height(DOC, BODY_NARROW, W9)
sh[28].width, sh[28].height = Inches(photo_w), Inches(photo_h)
sh[9].left, sh[9].width = Inches(x9), Inches(W9)

# --------------------------- EM section: current-density map beside the chart
EM_FIG_H = 3.80
EM_GAP = 0.40
curr = add_pic(FIGDIR + "icem_fig3_currentdensity.png",
               EM_FIG_H * (2036 / 1168))
curr_w = Emu(curr.width).inches
chart_w = COL1_W - EM_GAP - curr_w
sh[16].width, sh[16].height = Inches(round(chart_w, 4)), Inches(EM_FIG_H)
cap_curr = add_caption(CAP_CURRENT, curr_w)
cap_chart = add_caption(CAP_CHART, chart_w)
I_CURR = len(sh); sh.append(curr)
I_CAP_CURR = len(sh); sh.append(cap_curr)
I_CAP_CHART = len(sh); sh.append(cap_chart)

# --------------------------- temperature fields under the optimal-size table
temp = add_pic(FIGDIR + "ecce_fig9_tempmaps.png", COL2_W)
cap_temp = add_caption(CAP_TEMP, COL2_W)
I_TEMP = len(sh); sh.append(temp)
I_CAP_TEMP = len(sh); sh.append(cap_temp)

# ---------------------------------------------------------------- apply text
W7 = W15 = COL1_W
W18 = W25 = COL2_W
WH_L, WH_R = COL1_W, COL2_W

set_items(sh[7], MOTIVATION, BODY)
set_items(sh[9], DOC, BODY_NARROW)
set_items(sh[15], EM, BODY)
set_items(sh[18], OPT, BODY)
set_items(sh[25], CONCL, BODY)
for idx, txt in HEADERS.items():
    set_items(sh[idx], [txt], HEAD, bullets=False)
set_items(sh[13], ["Can the channel be shaped freely for cooling?"], QUESTION,
          bullets=False, algn="ctr")

# tables ---------------------------------------------------------------
t1 = sh[11].table
for r, (k, v) in enumerate([
    ("Power / speed", "300 kW / 2000 rpm"),
    ("Slots × layers", "48 × 4"),
    ("Conductor", "6 × 4 mm"),
    ("Current / frequency", "350 A, 266.7 Hz"),
    ("Coolant", "Oil, 65 °C, 10 L/min"),
    ("Flow regime", "Laminar"),
]):
    set_cell(t1, r, 0, k)
    set_cell(t1, r, 1, v)
set_table_font(t1, TABLE, 0.72)

t2 = sh[20].table
set_cell(t2, 0, 1, "Optimal size")
set_cell(t2, 0, 3, "Pumping")
set_table_font(t2, TABLE, 0.72)

# chart: shorter series names and a legend that fits the half-width frame.
# NOTE: edit the element tree - assigning part._blob is ignored for XML parts.
C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
SHORT_SERIES = {"Round channel": "Round", "Rectangular channel": "Rectangular"}
chartspace = sh[16].chart._chartSpace
for v in chartspace.iter(C + "v"):
    if v.text in SHORT_SERIES:
        v.text = SHORT_SERIES[v.text]
for defRPr in chartspace.iter(A + "defRPr"):
    defRPr.set("sz", str(CHART_FONT * 100))
# the rotated y-axis title collides with the tick labels at half width;
# drop it and carry the unit in the caption instead
for ax in chartspace.iter(C + "valAx"):
    for title in ax.findall(C + "title"):
        ax.remove(title)
    autodel = ax.find(C + "autoTitleDeleted")
    if autodel is not None:
        autodel.set("val", "1")


# ---------------------------------------------------------------- re-flow
def H(idx):
    return Emu(sh[idx].height).inches


def move(idx, top_in, dx=None):
    sh[idx].top = Inches(round(top_in, 4))
    if dx is not None:
        sh[idx].left = Inches(round(dx, 4))


h_head_l = text_height(["X"], HEAD, WH_L, bold=True, bullets=False)
h_head_doc = text_height([HEADERS[8]], HEAD, WH_L, bold=True, bullets=False)
h_head_r = text_height(["X"], HEAD, WH_R, bold=True, bullets=False)

h7 = text_height(MOTIVATION, BODY, W7)
h15 = text_height(EM, BODY, W15)
h18 = text_height(OPT, BODY, W18)
h25 = text_height(CONCL, BODY, W25)
h13 = text_height(["Can the channel be shaped freely for cooling?"], QUESTION,
                  COL1_W - 0.40, bold=True, italic=True, bullets=False)

GAP_HB = 0.28                       # header -> body
tile_q_h = h13 + 0.55
em_top = h_head_l + GAP_HB + h15 + 0.35   # top of the figure row
em_block = em_top + EM_FIG_H + 0.10 + H(I_CAP_CURR)

# each block: (height, [(shape_idx, offset_from_block_top, optional_new_left)])
col1 = [
    (h_head_l + GAP_HB + h7, [(6, 0), (7, h_head_l + GAP_HB)]),
    (h_head_doc + GAP_HB + photo_h,
     [(8, 0), (28, h_head_doc + GAP_HB),
      (9, h_head_doc + GAP_HB + max(0, (photo_h - h9) / 2))]),
    (h_head_l + GAP_HB + 6 * 0.72, [(10, 0), (11, h_head_l + GAP_HB)]),
    (em_block,
     [(14, 0), (15, h_head_l + GAP_HB),
      (I_CURR, em_top, COL1_X),                              # current density
      (I_CAP_CURR, em_top + EM_FIG_H + 0.10, COL1_X),        # its caption
      (16, em_top, COL1_X + curr_w + EM_GAP),                # AC-loss chart
      (I_CAP_CHART, em_top + EM_FIG_H + 0.10, COL1_X + curr_w + EM_GAP)]),
    (tile_q_h, [(12, 0), (13, (tile_q_h - h13) / 2)]),
]
col2 = [
    (h_head_r + GAP_HB + h18, [(17, 0), (18, h_head_r + GAP_HB)]),
    (H(29) + 0.45 + H(19), [(29, 0), (19, H(29) + 0.45)]),
    (4 * 0.72, [(20, 0)]),
    (H(I_TEMP) + 0.10 + H(I_CAP_TEMP),
     [(I_TEMP, 0, COL2_X), (I_CAP_TEMP, H(I_TEMP) + 0.10, COL2_X)]),
    (h_head_r + GAP_HB + h25, [(24, 0), (25, h_head_r + GAP_HB)]),
]

TOP, BOT = 10.55, 42.30


def flow(blocks, name):
    natural = sum(b[0] for b in blocks)
    slack = (BOT - TOP) - natural
    gap = max(0.55, min(1.30, slack / (len(blocks) - 1)))
    y = TOP
    for h, members in blocks:
        for m in members:
            move(m[0], y + m[1], m[2] if len(m) > 2 else None)
        y += h + gap
    end = y - gap
    flag = "  << OVERFLOW" if end > BOT + 0.01 else ""
    print(f"{name}: content {natural:.2f}in, gap {gap:.2f}in, ends {end:.2f}in{flag}")


flow(col1, "col1")
flow(col2, "col2")

sh[12].height = Inches(round(tile_q_h, 4))

prs.save(OUT)
print("saved", OUT)
